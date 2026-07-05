"""Этапы 0+1: извлечение текста из сырого корпуса → docs.meta.jsonl (git) + docs.text.jsonl (локально).
MAGIC-сниффинг типа, split-архивы .001/.002, дедуп по md5 нормализованного текста, year из имени.
Идемпотентно, параллельно. Профили корпуса потеряны → kg_value/sensitivity по категории (честный прокси).
"""
from __future__ import annotations
import os, sys, json, hashlib, zipfile, subprocess, tempfile, re, io
from concurrent.futures import ProcessPoolExecutor, as_completed
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from src.config import CORPUS_DIR, DOCS_META, DOCS_TEXT, nfc
from src.obs import get_logger

log = get_logger("extract")

DOC_EXT = {".pdf", ".docx", ".docm", ".doc", ".pptx", ".xls", ".xlsx", ".rtf"}

# --- Распознавание типов документов ТЗ (патенты/нормативы/справочники и пр.) ---
# Один словарь маркеров: тип → список регексов (RU+EN), первый совпавший тип и берётся.
# Порядок важен: специфичные (патент/норматив) раньше общих (статья/презентация).
DOC_TYPE_MARKERS = {
    "patent": [r"\bпатент\b", r"\bRU\s?\d{6,}", r"\bRU2\d{6}\b", r"\bМПК\b",
               r"заявка\s*№", r"приоритет\s+от", r"\(51\)\s*Int\.?\s*Cl", r"\bIPC\b"],
    "standard": [r"\bГОСТ\b", r"\bТУ\s?20", r"\bСП\s?\d+", r"\bСанПиН\b",
                 r"стандарт\s+организации", r"\bСТО\b", r"\bстандарт\b"],
    "reference": [r"\bсправочник\b", r"\bэнциклопед", r"каталог\s+оборудования", r"\bкаталог\b"],
    "protocol": [r"протокол\s+испытан", r"протокол\s+эксперимент", r"\bпротокол\b"],
    "presentation": [r"--\s*слайд\b", r"\bслайд\s*\d", r"\bпрезентац"],
    "report": [r"\bотч[её]т\b", r"\bдоклад\b", r"technical\s+report"],
    "article": [r"\bаннотац", r"\babstract\b", r"\bDOI\b", r"список\s+литератур",
                r"\bстатья\b", r"references\b"],
}
_DOC_TYPE_RE = {t: [re.compile(p, re.IGNORECASE) for p in pats]
                for t, pats in DOC_TYPE_MARKERS.items()}

# Метаданные патента: что найдётся, иначе None.
_PATENT_META_RE = {
    "number": re.compile(r"\b(RU\s?2?\d{6,})", re.IGNORECASE),
    "ipc": re.compile(r"(?:МПК|IPC|\(51\)\s*Int\.?\s*Cl\.?)[:\s]*([A-ZА-Я]\d{2}[A-ZА-Я]?\s?\d{0,3}/?\d{0,4})",
                      re.IGNORECASE),
    "applicant": re.compile(r"(?:заявитель|патентообладатель|applicant)[:\s]*([^\n]{3,120})",
                            re.IGNORECASE),
    "priority_date": re.compile(r"приоритет\s+от[:\s]*([\d.\-/]{6,10})", re.IGNORECASE),
}


def classify_doc_type(text: str, filename: str = "") -> str:
    """Тип документа по маркерам в тексте+имени файла.
    Возвращает один из {patent,standard,reference,report,article,presentation,protocol,other}."""
    hay = f"{filename}\n{text or ''}"
    for dtype, regexes in _DOC_TYPE_RE.items():
        if any(rx.search(hay) for rx in regexes):
            return dtype
    return "other"


def extract_patent_meta(text: str) -> dict:
    """Метаданные патента по регексам: {number, ipc, applicant, priority_date} (None если не найдено)."""
    t = text or ""
    out = {}
    for key, rx in _PATENT_META_RE.items():
        m = rx.search(t)
        out[key] = (m.group(1).strip() if m else None)
    return out
MAGIC = {b"%PDF": "pdf", b"PK\x03\x04": "zip", b"\xd0\xcf\x11\xe0": "ole2",
         b"Rar!": "rar", b"{\\rtf": "rtf"}
# прокси метаданных по категории (профили корпуса утеряны — восстановимо ре-профилированием)
CAT_KG = {"Статьи": 4, "Обзоры": 4, "Доклады": 3, "Журналы": 2, "Материалы конференций": 1}
CAT_SENS = {"Статьи": "internal", "Доклады": "internal", "Обзоры": "internal",
            "Журналы": "public", "Материалы конференций": "public"}

def sniff(path: str) -> str:
    try:
        with open(path, "rb") as f:
            head = f.read(8)
    except OSError:
        return ""
    for magic, kind in MAGIC.items():
        if head.startswith(magic):
            return kind
    return os.path.splitext(path)[1].lower().lstrip(".")

def year_from_name(name: str):
    years = re.findall(r"(?:19|20)\d{2}", name)
    return int(years[-1]) if years else None

def _extract_pdf(p):
    import fitz
    d = fitz.open(p); return "\n".join(pg.get_text() for pg in d), len(d)
def _extract_docx(p):
    import docx
    doc = docx.Document(p); parts = [x.text for x in doc.paragraphs]
    for t in doc.tables:
        for r in t.rows: parts.append(" | ".join(c.text for c in r.cells))
    return "\n".join(parts), 0
def _extract_pptx(p):
    from pptx import Presentation
    prs = Presentation(p); out = []
    for i, s in enumerate(prs.slides, 1):
        out.append(f"--слайд {i}--")
        for sh in s.shapes:
            if sh.has_text_frame: out.append(sh.text_frame.text)
            if getattr(sh, "has_table", False):
                for r in sh.table.rows: out.append(" | ".join(c.text for c in r.cells))
    return "\n".join(out), len(prs.slides)
def _extract_xls(p):
    import pandas as pd
    xl = pd.ExcelFile(p); out = []
    for name in xl.sheet_names[:8]:
        try: out.append(f"=== лист {name} ===\n" + xl.parse(name, nrows=60).to_string(max_rows=60))
        except Exception as e: out.append(f"[лист {name}: {e}]")
    return "\n".join(out), len(xl.sheet_names)
def _extract_ole_or_rtf(p):
    with tempfile.TemporaryDirectory() as td:
        subprocess.run(["soffice", "--headless", f"-env:UserInstallation=file://{td}/lo",
                        "--convert-to", "docx", "--outdir", td, p],
                       capture_output=True, timeout=180)
        out = os.path.join(td, os.path.splitext(os.path.basename(p))[0] + ".docx")
        if os.path.exists(out): return _extract_docx(out)
    return "", 0

def extract_any(path: str):
    """Диспетчер по MAGIC (расширение — только подсказка)."""
    kind = sniff(path)
    ext = os.path.splitext(path)[1].lower()
    try:
        if kind == "pdf": return _extract_pdf(path)
        if kind == "zip":  # docx/xlsx/pptx — все zip-based; различаем по расширению
            if ext in (".docx", ".docm"): return _extract_docx(path)
            if ext == ".pptx": return _extract_pptx(path)
            if ext in (".xlsx",): return _extract_xls(path)
            return _extract_docx(path)
        if kind == "ole2":
            if ext in (".xls",): return _extract_xls(path)
            return _extract_ole_or_rtf(path)
        if kind == "rtf": return _extract_ole_or_rtf(path)
        # fallback по расширению
        if ext in (".docx", ".docm"): return _extract_docx(path)
        if ext == ".pptx": return _extract_pptx(path)
        if ext in (".xls", ".xlsx"): return _extract_xls(path)
        if ext in (".doc", ".rtf"): return _extract_ole_or_rtf(path)
    except Exception as e:
        return f"[EXTRACT-ERR {type(e).__name__}: {e}]", 0
    return "", 0

def collect_split_archives(root):
    """Склеить .001/.002 → zip → распаковать во временные doc-файлы. Возвращает (disp, realpath, cat)."""
    jobs, tmpdirs = [], []
    seen001 = {}
    for dp, _, fs in os.walk(root):
        for fn in fs:
            if fn.lower().endswith(".001"):
                seen001[os.path.join(dp, fn)] = True
    for first in seen001:
        base = first[:-4]
        parts = sorted(p for p in [base + f".{i:03d}" for i in range(1, 20)] if os.path.exists(p))
        if not parts: continue
        td = tempfile.mkdtemp(prefix="split_")
        tmpdirs.append(td)
        zippath = os.path.join(td, "joined.zip")
        with open(zippath, "wb") as out:
            for part in parts:
                with open(part, "rb") as f: out.write(f.read())
        try:
            with zipfile.ZipFile(zippath) as z: z.extractall(td)
        except Exception:
            continue
        cat = os.path.relpath(first, root).split(os.sep)[0]
        for adp, _, afs in os.walk(td):
            for afn in afs:
                if os.path.splitext(afn)[1].lower() in DOC_EXT:
                    disp = os.path.relpath(base, root) + " :: " + afn
                    jobs.append((disp, os.path.join(adp, afn), cat))
    return jobs

def process(job):
    disp, real, cat = job
    try:
        text, pages = extract_any(real)
    except Exception as e:  # ТЗ NFR: битый файл не роняет весь прогон
        log.warning("extract failed file=%s err=%s: %s", disp, type(e).__name__, e)
        text, pages = f"[EXTRACT-ERR {type(e).__name__}: {e}]", 0
    text = nfc(text)
    cyr = sum(1 for c in text[:4000] if "а" <= c.lower() <= "я")
    lat = sum(1 for c in text[:4000] if "a" <= c.lower() <= "z")
    doc_type = classify_doc_type(text, disp)
    out = {
        "src": nfc(disp), "cat": nfc(cat), "pages": pages, "chars": len(text),
        "lang": "RU" if cyr >= lat else "EN",
        "year": year_from_name(disp),
        "sensitivity": CAT_SENS.get(cat, "public"),
        "kg_value": CAT_KG.get(cat, 1),
        "doc_type": doc_type,
        "text": text,
        "ok": not text.startswith("[EXTRACT-ERR"),
    }
    if doc_type == "patent":
        out["patent_meta"] = extract_patent_meta(text)
    return out

def main(limit=0):
    root = str(CORPUS_DIR)
    direct = []
    for dp, _, fs in os.walk(root):
        for fn in fs:
            if os.path.splitext(fn)[1].lower() in DOC_EXT:
                rp = os.path.relpath(os.path.join(dp, fn), root)
                direct.append((rp, os.path.join(dp, fn), rp.split(os.sep)[0]))
    split = collect_split_archives(root)
    jobs = direct + split
    if limit: jobs = jobs[:limit]
    print(f"Документов: прямых {len(direct)}, из split-архивов {len(split)}, всего {len(jobs)}", flush=True)

    seen_md5, meta_rows, text_rows, done = {}, [], [], 0
    with ProcessPoolExecutor(max_workers=max(2, (os.cpu_count() or 4) - 1)) as ex:
        futs = [ex.submit(process, j) for j in jobs]
        for fu in as_completed(futs):
            r = fu.result(); done += 1
            norm = re.sub(r"\s+", " ", r["text"]).strip().lower()
            md5 = hashlib.md5(norm.encode()).hexdigest()
            doc_id = md5[:16]
            if md5 in seen_md5:  # дедуп: канон уже есть
                continue
            seen_md5[md5] = doc_id
            meta = {k: r[k] for k in ("src", "cat", "pages", "chars", "lang", "year",
                                       "sensitivity", "kg_value", "doc_type", "ok")}
            if "patent_meta" in r:
                meta["patent_meta"] = r["patent_meta"]
            meta["doc_id"] = doc_id
            meta_rows.append(meta)
            text_rows.append({"doc_id": doc_id, "text": r["text"]})
            if done % 200 == 0: print(f"  {done}/{len(jobs)}", flush=True)

    with open(DOCS_META, "w", encoding="utf-8") as f:
        for m in meta_rows: f.write(json.dumps(m, ensure_ascii=False) + "\n")
    with open(DOCS_TEXT, "w", encoding="utf-8") as f:
        for t in text_rows: f.write(json.dumps(t, ensure_ascii=False) + "\n")
    ok = sum(1 for m in meta_rows if m["ok"])
    dups = len(jobs) - len(meta_rows) if not limit else "?"
    print(f"ГОТОВО: уникальных {len(meta_rows)} (дублей отброшено {dups}), извлечено ок {ok}", flush=True)
    print(f"  meta → {DOCS_META}\n  text → {DOCS_TEXT}", flush=True)
    assert len(meta_rows) > 0 and ok > len(meta_rows) * 0.5, "слишком много ошибок извлечения"

def _selfcheck():
    """Self-check распознавания типов (без корпуса): 5 образцов + patent_meta."""
    samples = [
        ("Патент RU2123456 C1, МПК C22B 3/04, приоритет от 12.05.2019", "doc.pdf", "patent"),
        ("ГОСТ 12.1.005-88 Общие санитарно-гигиенические требования", "gost.docx", "standard"),
        ("Справочник металлурга. Энциклопедия цветных металлов", "spravochnik.pdf", "reference"),
        ("Аннотация. В статье рассмотрено... DOI:10.1/x. Список литературы", "art.pdf", "article"),
        ("--слайд 1--\nПрезентация результатов", "pres.pptx", "presentation"),
    ]
    for text, fn, expect in samples:
        got = classify_doc_type(text, fn)
        assert got == expect, f"classify({fn}): ждали {expect}, получили {got}"
    meta = extract_patent_meta("Патент RU2123456 C1")
    assert meta["number"] and "2123456" in meta["number"], f"номер патента не извлечён: {meta}"
    print("OK: classify_doc_type (5 образцов) + extract_patent_meta (номер RU2123456)")


if __name__ == "__main__":
    if len(sys.argv) > 1 and sys.argv[1] == "selfcheck":
        _selfcheck()
    else:
        lim = int(sys.argv[1]) if len(sys.argv) > 1 else 0
        main(lim)
